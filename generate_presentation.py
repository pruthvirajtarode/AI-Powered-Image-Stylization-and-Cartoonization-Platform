"""
Toonify AI - Professional PowerPoint Presentation Generator
Generates a premium, modern presentation with Gamma-style design
Run: python generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Create presentation object
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(15, 23, 42)  # Dark navy
ACCENT_BLUE = RGBColor(59, 130, 246)  # Bright blue
ACCENT_PURPLE = RGBColor(147, 51, 234)  # Purple
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(229, 231, 235)
GOLD = RGBColor(251, 191, 36)

def add_gradient_background(slide, color1, color2):
    """Add gradient background to slide"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 45.0
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_gradient_background(slide, DARK_BG, RGBColor(30, 58, 138))
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "Toonify AI"
    p.font.size = Pt(96)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "Transforming Reality into Art with AI"
    p.font.size = Pt(32)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Add visual accent line
    line = slide.shapes.add_shape(1, Inches(3.5), Inches(5.5), Inches(3), Inches(0))
    line.line.color.rgb = GOLD
    line.line.width = Pt(3)

def add_problem_slide(prs):
    """Slide 2: Problem Statement"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, DARK_BG, RGBColor(25, 55, 119))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "The Problem"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    # Problem points
    problems = [
        "❌ Complex tools like Photoshop—steep learning curve",
        "❌ Expensive subscriptions ($50-500/month)",
        "❌ Time-consuming editing process (hours)",
        "❌ Desktop-only, limited accessibility"
    ]
    
    y_pos = 1.8
    for problem in problems:
        text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8.5), Inches(0.6))
        text_frame = text_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = problem
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        y_pos += 0.9

def add_overview_slide(prs):
    """Slide 3: Project Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, DARK_BG, RGBColor(30, 58, 138))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "The Solution"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    
    # Main description
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.7))
    desc_frame = desc_box.text_frame
    p = desc_frame.paragraphs[0]
    p.text = "Web-based AI Image Stylization Platform"
    p.font.size = Pt(28)
    p.font.color.rgb = GOLD
    
    # Workflow
    workflow_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(9), Inches(0.6))
    workflow_frame = workflow_box.text_frame
    p = workflow_frame.paragraphs[0]
    p.text = "Upload → Select Style → Generate → Download"
    p.font.size = Pt(22)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Key features
    features_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(0.5))
    features_frame = features_title.text_frame
    p = features_frame.paragraphs[0]
    p.text = "Key Features:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    features = [
        "✨ 9+ AI Styles (Pixar, Anime, Sketch, Oil Paint, etc.)",
        "⚡ Ultra-fast processing (3-4 seconds per image)",
        "🎨 High-quality 4K outputs with facial preservation"
    ]
    
    y_pos = 3.6
    for feature in features:
        text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8.5), Inches(0.5))
        text_frame = text_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = feature
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        y_pos += 0.7

def add_workflow_slide(prs):
    """Slide 4: Demo Workflow"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, DARK_BG, RGBColor(30, 58, 138))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "How It Works"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    # Workflow steps in boxes
    steps = [
        ("📸", "Upload\nImage"),
        ("🎨", "Select\nAI Style"),
        ("⚙️", "Process\nwith AI"),
        ("⬇️", "Download\n4K Output")
    ]
    
    x_start = 1
    box_width = 1.8
    box_height = 2.5
    
    for i, (emoji, text) in enumerate(steps):
        x_pos = x_start + (i * 2.1)
        
        # Box
        box = slide.shapes.add_shape(1, Inches(x_pos), Inches(1.8), Inches(box_width), Inches(box_height))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(55, 65, 81)
        box.line.color.rgb = ACCENT_PURPLE
        box.line.width = Pt(2)
        
        # Emoji
        emoji_box = slide.shapes.add_textbox(Inches(x_pos), Inches(2.2), Inches(box_width), Inches(0.8))
        emoji_frame = emoji_box.text_frame
        p = emoji_frame.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(48)
        p.alignment = PP_ALIGN.CENTER
        
        # Text
        text_box = slide.shapes.add_textbox(Inches(x_pos), Inches(3.2), Inches(box_width), Inches(1))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Arrow (except last)
        if i < len(steps) - 1:
            arrow_x = x_pos + box_width + 0.05
            arrow = slide.shapes.add_shape(1, Inches(arrow_x), Inches(3.3), Inches(0.15), Inches(0))
            arrow.line.color.rgb = GOLD
            arrow.line.width = Pt(2)

def add_results_slide(prs):
    """Slide 5: Output Results"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, DARK_BG, RGBColor(30, 58, 138))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Stunning Results"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = GOLD
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "Before & After Transformations"
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Image placeholder boxes
    placeholders = [
        ("Original", 1.2),
        ("Pixar 3D", 3.8),
        ("Anime", 6.4)
    ]
    
    for label, x_pos in placeholders:
        # Box
        box = slide.shapes.add_shape(1, Inches(x_pos), Inches(2.0), Inches(2.5), Inches(3.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(75, 85, 99)
        box.line.color.rgb = ACCENT_BLUE
        box.line.width = Pt(3)
        
        # Label
        label_box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.7), Inches(2.5), Inches(0.4))
        label_frame = label_box.text_frame
        p = label_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT_PURPLE
        p.alignment = PP_ALIGN.CENTER
        
        # Placeholder text
        text_box = slide.shapes.add_textbox(Inches(x_pos + 0.3), Inches(3.5), Inches(1.9), Inches(1))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = "[Image Placeholder]"
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    # Highlights
    highlights = [
        "✓ 95% Facial Feature Preservation",
        "✓ Processing Speed: 3.2 seconds",
        "✓ 4K Professional Quality"
    ]
    
    y_pos = 5.5
    for highlight in highlights:
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.4))
        text_frame = text_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = highlight
        p.font.size = Pt(18)
        p.font.color.rgb = GOLD
        y_pos += 0.5

def add_tech_stack_slide(prs):
    """Slide 6: Technology Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, DARK_BG, RGBColor(30, 58, 138))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Technology Stack"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    # Tech sections
    tech_sections = [
        {
            "title": "🖥️ Frontend",
            "items": ["HTML5 • CSS3", "JavaScript • Canvas API"],
            "x": 0.5, "y": 1.2
        },
        {
            "title": "⚙️ Backend",
            "items": ["Python 3.9+", "Flask • Gunicorn"],
            "x": 3.3, "y": 1.2
        },
        {
            "title": "🤖 AI Processing",
            "items": ["OpenCV 4.10", "NumPy • Pillow"],
            "x": 6.1, "y": 1.2
        },
        {
            "title": "🗄️ Database",
            "items": ["PostgreSQL", "SQLite (Backup)"],
            "x": 0.5, "y": 3.8
        },
        {
            "title": "🔐 Security",
            "items": ["OAuth 2.0", "Bcrypt • SSL/TLS"],
            "x": 3.3, "y": 3.8
        },
        {
            "title": "🚀 Deployment",
            "items": ["Docker", "Render • CDN"],
            "x": 6.1, "y": 3.8
        }
    ]
    
    for section in tech_sections:
        # Box
        box = slide.shapes.add_shape(1, Inches(section["x"]), Inches(section["y"]), Inches(2.5), Inches(1.95))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(55, 65, 81)
        box.line.color.rgb = ACCENT_PURPLE
        box.line.width = Pt(2)
        
        # Title
        title = slide.shapes.add_textbox(Inches(section["x"] + 0.15), Inches(section["y"] + 0.15), Inches(2.2), Inches(0.45))
        title_frame = title.text_frame
        p = title_frame.paragraphs[0]
        p.text = section["title"]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        
        # Items
        y_offset = section["y"] + 0.7
        for item in section["items"]:
            item_box = slide.shapes.add_textbox(Inches(section["x"] + 0.15), Inches(y_offset), Inches(2.2), Inches(0.35))
            item_frame = item_box.text_frame
            p = item_frame.paragraphs[0]
            p.text = item
            p.font.size = Pt(13)
            p.font.color.rgb = LIGHT_GRAY
            y_offset += 0.4

def add_challenges_slide(prs):
    """Slide 7: Challenges Faced"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, DARK_BG, RGBColor(25, 55, 119))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Challenges Faced"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(239, 68, 68)
    
    challenges = [
        "⚡ Handling diverse image qualities and resolutions",
        "🎯 Maintaining 95%+ facial feature accuracy",
        "⏱️ Optimizing processing speed to 3-4 seconds",
        "🎨 Designing intuitive UI for non-technical users"
    ]
    
    y_pos = 1.6
    for challenge in challenges:
        text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8.5), Inches(0.7))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = challenge
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        y_pos += 1.1

def add_solutions_slide(prs):
    """Slide 8: Solutions Implemented"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, DARK_BG, RGBColor(30, 58, 138))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Solutions Implemented"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(34, 197, 94)
    
    solutions = [
        "✓ Developed and trained custom deep learning models",
        "✓ Implemented efficient multi-threaded backend pipeline",
        "✓ Applied advanced image preprocessing techniques",
        "✓ Built clean, intuitive, and responsive UI"
    ]
    
    y_pos = 1.6
    for solution in solutions:
        text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8.5), Inches(0.7))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = solution
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        y_pos += 1.1

def add_future_scope_slide(prs):
    """Slide 9: Future Scope"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, DARK_BG, RGBColor(30, 58, 138))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Future Roadmap"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    
    features = [
        "🎬 Real-time video stylization",
        "📱 Native mobile applications (iOS & Android)",
        "🚀 Advanced batch processing for enterprises",
        "🌍 Extended AI style library (100+ models)"
    ]
    
    y_pos = 1.6
    for feature in features:
        text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8.5), Inches(0.7))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = feature
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        y_pos += 1.1

def add_conclusion_slide(prs):
    """Slide 10: Conclusion"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, DARK_BG, RGBColor(30, 58, 138))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "AI is democratizing creativity"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = "Toonify AI transforms complex image editing into a simple, fast, and accessible experience for creators, designers, and businesses worldwide."
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_thankyou_slide(prs):
    """Slide 11: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, ACCENT_PURPLE, ACCENT_BLUE)
    
    # Main text
    thank_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    thank_frame = thank_box.text_frame
    p = thank_frame.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(96)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Q&A
    qa_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    qa_frame = qa_box.text_frame
    p = qa_frame.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(56)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

def generate_presentation():
    """Generate the complete presentation"""
    print("🎨 Generating Toonify AI Presentation...")
    
    add_title_slide(prs)
    print("✓ Slide 1: Title")
    
    add_problem_slide(prs)
    print("✓ Slide 2: Problem")
    
    add_overview_slide(prs)
    print("✓ Slide 3: Overview")
    
    add_workflow_slide(prs)
    print("✓ Slide 4: Workflow")
    
    add_results_slide(prs)
    print("✓ Slide 5: Results")
    
    add_tech_stack_slide(prs)
    print("✓ Slide 6: Technology Stack")
    
    add_challenges_slide(prs)
    print("✓ Slide 7: Challenges")
    
    add_solutions_slide(prs)
    print("✓ Slide 8: Solutions")
    
    add_future_scope_slide(prs)
    print("✓ Slide 9: Future Scope")
    
    add_conclusion_slide(prs)
    print("✓ Slide 10: Conclusion")
    
    add_thankyou_slide(prs)
    print("✓ Slide 11: Thank You")
    
    # Save presentation
    output_path = "Toonify_AI_Presentation.pptx"
    prs.save(output_path)
    print(f"\n✅ Presentation saved: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print("\n🚀 Your professional presentation is ready!")

if __name__ == "__main__":
    generate_presentation()
